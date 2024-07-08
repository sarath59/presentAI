import os
import io
from flask import Flask, render_template, request, send_file, Response
from crewai import Agent, Task, Crew, Process
from langchain.tools import Tool
from openai import OpenAI
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import queue

app = Flask(__name__)

# Set up OpenAI API
client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

# Color scheme
color_scheme = {
    'primary': '#0070C0',    # Blue
    'secondary': '#FFC000',  # Gold
    'text': '#404040',       # Dark Gray
    'background': '#F8F8F8'  # Light Gray
}

# Queue for terminal output
terminal_output = queue.Queue()

# Helper function to convert hex to RGB
def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

# Web scraping tool (simplified for this example)
def scrape_website(url):
    terminal_output.put(f"Scraping website: {url}")
    return f"Content scraped from {url}"

scrape_tool = Tool(
    name="Scrape Website",
    func=scrape_website,
    description="Scrape content from a given URL"
)

# Basic web search tool (simplified for this example)
def basic_web_search(query):
    terminal_output.put(f"Searching for: {query}")
    return f"Search results for: {query}"

search_tool = Tool(
    name="Web Search",
    func=basic_web_search,
    description="Search for recent information on the internet"
)

# Enhanced OpenAI content generation tool
def generate_content(prompt):
    terminal_output.put("Generating content with OpenAI")
    enhanced_prompt = (
        f"{prompt}\n\n"
        "Create engaging and professional content for a presentation slide that includes:\n"
        "- A catchy title with an appropriate emoji\n"
        "- 3-5 key points, each as a bullet point with a brief explanation\n"
        "- Use professional language and tone\n"
        "- Include relevant statistics or data if applicable\n"
        "- Use emojis where appropriate to enhance engagement\n"
        "- End with a thought-provoking question or call-to-action\n"
        "Format the content for easy reading in a presentation slide, using bullet points and clear structure."
    )
    try:
        response = client.chat.completions.create(
            messages=[
                {
                    "role": "user",
                    "content": enhanced_prompt,
                }
            ],
            model="gpt-3.5-turbo",
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        error_message = f"Error generating content: {str(e)}"
        terminal_output.put(error_message)
        return error_message

content_tool = Tool(
    name="Generate Content",
    func=generate_content,
    description="Generate engaging and structured content for presentation slides"
)

# Presentation creation functions
def create_title_slide(prs, title, subtitle):
    terminal_output.put("Creating title slide")
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = title
    rgb = hex_to_rgb(color_scheme['primary'])
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(*rgb)
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    
    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = subtitle
    rgb = hex_to_rgb(color_scheme['secondary'])
    subtitle_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(*rgb)
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(32)

def create_content_slide(prs, title, content):
    terminal_output.put(f"Creating content slide: {title}")
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = title
    rgb = hex_to_rgb(color_scheme['primary'])
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(*rgb)
    title_shape.text_frame.paragraphs[0].font.size = Pt(40)
    
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()  # Clear existing text
    
    lines = content.split('\n')
    for line in lines:
        p = tf.add_paragraph()
        p.text = line.strip()
        p.font.size = Pt(24)
        rgb = hex_to_rgb(color_scheme['text'])
        p.font.color.rgb = RGBColor(*rgb)
        
        if line.startswith('•'):
            p.level = 1
        elif line.startswith('-'):
            p.level = 2
        else:
            p.level = 0
            p.font.bold = True
    
    return slide

# Define CrewAI agents
researcher = Agent(
    role="Information Researcher",
    goal="Find recent and relevant information for the presentation",
    backstory="You are an expert at finding the most up-to-date and relevant information on any topic.",
    verbose=True,
    allow_delegation=False,
    tools=[scrape_tool, search_tool]
)

content_creator = Agent(
    role="Content Creator",
    goal="Create engaging and informative content for each slide",
    backstory="You are a master of creating compelling content that captivates audiences.",
    verbose=True,
    allow_delegation=False,
    tools=[content_tool]
)

def generate_presentation(agenda):
    terminal_output.put(f"Generating presentation for agenda: {agenda}")
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    slide_types = [
        "Introduction",
        "Agenda Overview",
        "Key Point 1",
        "Key Point 2",
        "Key Point 3",
        "Supporting Data",
        "Case Study",
        "Conclusion",
        "Next Steps",
        "Q&A"
    ]

    crew = Crew(
        agents=[researcher, content_creator],
        tasks=[
            Task(
                description=f"Research recent information about: {agenda}",
                expected_output="A comprehensive report on the latest findings related to the agenda topic.",
                agent=researcher
            ),
            Task(
                description="Create engaging content for each slide based on the research",
                expected_output="Structured and engaging slide content for each slide type.",
                agent=content_creator
            )
        ],
        process=Process.sequential
    )

    result = crew.kickoff(inputs={'agenda': agenda})

    # Create title slide
    create_title_slide(prs, f"📊 {agenda}", "An AI-Generated Presentation")

    # Parse the result and create content slides
    slides_content = result.split("\n\n")
    for i, content in enumerate(slides_content):
        if i < len(slide_types):
            slide_type = slide_types[i]
            slide = create_content_slide(prs, slide_type, content)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    terminal_output.put("Presentation generated successfully")
    return output

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/generate', methods=['POST'])
def generate():
    agenda = request.form['agenda']
    presentation = generate_presentation(agenda)
    return send_file(
        presentation,
        as_attachment=True,
        download_name='AI_Generated_Presentation.pptx',
        mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
    )

@app.route('/stream')
def stream():
    def generate():
        while True:
            message = terminal_output.get()
            yield f"data: {message}\n\n"
    return Response(generate(), mimetype='text/event-stream')

if __name__ == '__main__':
    app.run(debug=True, threaded=True)